#!/usr/bin/env python
"""
Script to populate PPTX template with candidate information while preserving formatting.
This script reads the Prenume_Nume_Technical_Profile.pptx template and fills it with candidate data.
"""

from pptx import Presentation
from pptx.util import Pt
from pathlib import Path
import shutil
import time

def extract_candidate_info():
    """Extract candidate information from the converted markdown file."""
    
    candidate_data = {
        'name': 'VLAD GHEORGHE BRANZEI',
        'title': 'Senior .NET Full Stack Developer',
        'overview': 'Senior Full Stack Developer with 13+ years of enterprise experience in microservices, fintech, and stack modernization. Specialized in transforming legacy monolithic systems into cloud-native, microservices-based solutions.',
        'core_skills': [
            'C# / .NET 6/8/10 / .NET Framework',
            'Entity Framework / LINQ',
            'Angular 16+ / 19+',
            'SQL Server / T-SQL',
            'Microservices Architecture'
        ],
        'additional_skills': [
            'Azure Cloud (AKS, ACR, DevOps)',
            'Docker / Kubernetes',
            'RabbitMQ / Kafka',
            'CI/CD (Jenkins, Github Actions)',
            'Python (FastAPI, LangGraph)'
        ],
        'soft_skills': [
            'Architectural Design & Leadership',
            'Technical Mentoring & Code Review',
            'Problem Solving & Critical Thinking',
            'Continuous Learning & Innovation',
            'Agile Methodology'
        ],
        'languages': [
            'English – Professional',
            'German – Beginner',
            'Romanian – Native'
        ],
        'education': 'Technical University of Cluj-Napoca',
        'certifications': 'OWASP 10 Web Security | Azure Fundamentals',
        'experience': [
            {
                'company': 'SOFTLAB360',
                'title': 'Senior .NET Full Stack Developer',
                'period': 'Jan 2024 – Present',
                'area': 'Fintech / Wealth Management',
                'achievements': [
                    'Monolith refactorization (WebForms → REST + Angular)',
                    'REST API development with architectural guidelines',
                    'Kafka integration (topics, publishing, consuming, Avro)',
                    'Python scripts for automated API testing',
                    'Multitenancy implementation',
                    'Critical issue resolution (thread starvation)'
                ]
            },
            {
                'company': 'PROGRAMMING POOL, INC.',
                'title': 'Senior Full Stack .NET Developer',
                'period': 'Apr 2020 – Jan 2024',
                'area': 'Employee Benefits Business',
                'achievements': [
                    'Monolith → Microservices architecture transition',
                    'UI refactoring (EpiServer → Angular 19)',
                    'MasterCard Card issuer integration',
                    'RabbitMQ with MassTransit implementation',
                    'OWASP vulnerability scanning and remediation',
                    'Code reviews and technical interviews'
                ]
            },
            {
                'company': 'MICROFOCUS',
                'title': 'Senior SQL Developer',
                'period': 'Sept 2019 – Apr 2020',
                'area': 'UCMDB Product Support',
                'achievements': [
                    'Database performance optimization',
                    'High availability configurations (SQL Server AlwaysOn, PostgreSQL Patroni)',
                    'Performance analysis and tuning',
                    'Diagnostics script creation',
                    'Deadlock and timeout resolution'
                ]
            }
        ]
    }
    
    return candidate_data

def update_text_frame(text_frame, new_text, preserve_format=False):
    """Update text frame while optionally preserving formatting."""
    # Preserve formatting from first paragraph
    if preserve_format and len(text_frame.paragraphs) > 0:
        first_para = text_frame.paragraphs[0]
        font_size = first_para.font.size if first_para.font.size else Pt(11)
        font_bold = first_para.font.bold if first_para.font.bold is not None else False
        font_color = first_para.font.color if first_para.font.color else None
    
    # Clear existing text
    text_frame.clear()
    
    # Add new text
    p = text_frame.paragraphs[0]
    p.text = new_text
    
    if preserve_format:
        if font_size:
            p.font.size = font_size
        if font_bold is not None:
            p.font.bold = font_bold

def copy_formatting(source_run, target_run):
    """Copy formatting properties from source run to target run."""
    # Font properties
    if source_run.font.name:
        target_run.font.name = source_run.font.name
    if source_run.font.size:
        target_run.font.size = source_run.font.size
    if source_run.font.bold is not None:
        target_run.font.bold = source_run.font.bold
    if source_run.font.italic is not None:
        target_run.font.italic = source_run.font.italic
    if source_run.font.underline is not None:
        target_run.font.underline = source_run.font.underline
    
    # Color
    if source_run.font.color and source_run.font.color.rgb:
        target_run.font.color.rgb = source_run.font.color.rgb

def copy_paragraph_formatting(source_para, target_para):
    """Copy paragraph-level formatting properties."""
    # Alignment
    if source_para.alignment:
        target_para.alignment = source_para.alignment
    
    # Level
    if source_para.level is not None:
        target_para.level = source_para.level
    
    # Space before/after
    if source_para.space_before:
        target_para.space_before = source_para.space_before
    if source_para.space_after:
        target_para.space_after = source_para.space_after

def update_paragraph_text(paragraph, new_text):
    """Update paragraph text while preserving formatting."""
    # Save formatting info from existing runs
    formatting_info = {}
    if len(paragraph.runs) > 0:
        source_run = paragraph.runs[0]
        formatting_info['font_name'] = source_run.font.name
        formatting_info['font_size'] = source_run.font.size
        formatting_info['bold'] = source_run.font.bold
        formatting_info['italic'] = source_run.font.italic
        formatting_info['underline'] = source_run.font.underline
        
        # Try to get color, handling both RGB and scheme colors
        try:
            if source_run.font.color.type:
                formatting_info['color_rgb'] = source_run.font.color.rgb
        except (AttributeError, ValueError):
            # Color is a scheme color or not available, skip it
            pass
    
    # Clear existing runs and set new text
    paragraph.clear()
    new_run = paragraph.add_run()
    new_run.text = new_text
    
    # Apply saved formatting to new run
    if formatting_info:
        if 'font_name' in formatting_info and formatting_info['font_name']:
            new_run.font.name = formatting_info['font_name']
        if 'font_size' in formatting_info and formatting_info['font_size']:
            new_run.font.size = formatting_info['font_size']
        if 'bold' in formatting_info and formatting_info['bold'] is not None:
            new_run.font.bold = formatting_info['bold']
        if 'italic' in formatting_info and formatting_info['italic'] is not None:
            new_run.font.italic = formatting_info['italic']
        if 'underline' in formatting_info and formatting_info['underline'] is not None:
            new_run.font.underline = formatting_info['underline']
        if 'color_rgb' in formatting_info:
            try:
                new_run.font.color.rgb = formatting_info['color_rgb']
            except (AttributeError, ValueError):
                # Color assignment failed, continue without color
                pass

def populate_template():
    """Load template and populate with candidate information."""
    template_path = r"c:\Users\turca\OneDrive\github\recruitment-ai-assistant\templates\Prenume_Nume_Technical_Profile.pptx"
    output_path = r"c:\Users\turca\OneDrive\github\recruitment-ai-assistant\output\Alupei Gheorghe\Alupei Gheorghe.pptx"
    temp_output = output_path + ".tmp"
    
    # Load template
    prs = Presentation(template_path)
    candidate = extract_candidate_info()
    
    print(f"Total slides: {len(prs.slides)}\n")
    
    # Process Slide 1 (Profile and skills)
    if len(prs.slides) > 0:
        slide1 = prs.slides[0]
        print("=== SLIDE 1 ===")
        
        # Find the main textbox (Shape 1)
        if len(slide1.shapes) > 1:
            textbox = slide1.shapes[1]
            if hasattr(textbox, "text_frame"):
                tf = textbox.text_frame
                
                # Update specific paragraphs by index
                para_updates = {
                    0: ("PRENUME NUME", candidate['name']),
                    1: ("Software Engineer", candidate['title']),
                    5: ("I'm an intuitive", candidate['overview']),
                    10: ("C++", candidate['core_skills'][0]),
                    11: ("C#", candidate['core_skills'][1] if len(candidate['core_skills']) > 1 else ""),
                    12: ("Low-level", candidate['core_skills'][2] if len(candidate['core_skills']) > 2 else ""),
                    16: ("Electronics", candidate['additional_skills'][0]),
                    17: ("Matlab", candidate['additional_skills'][1] if len(candidate['additional_skills']) > 1 else ""),
                    18: ("Linux", candidate['additional_skills'][2] if len(candidate['additional_skills']) > 2 else ""),
                    19: ("Vector CAN", candidate['additional_skills'][3] if len(candidate['additional_skills']) > 3 else ""),
                    24: ("Positive attitude", candidate['soft_skills'][0]),
                    25: ("Communicative", candidate['soft_skills'][1] if len(candidate['soft_skills']) > 1 else ""),
                    26: ("Continuous learner", candidate['soft_skills'][2] if len(candidate['soft_skills']) > 2 else ""),
                    31: ("English", candidate['languages'][0]),
                    32: ("French", candidate['languages'][1] if len(candidate['languages']) > 1 else ""),
                    37: ("Politehnica", candidate['education']),
                }
                
                for para_idx, (old_key, new_text) in para_updates.items():
                    if para_idx < len(tf.paragraphs):
                        para = tf.paragraphs[para_idx]
                        if old_key in para.text or old_key.lower() in para.text.lower():
                            update_paragraph_text(para, new_text)
                            print(f"  Updated para {para_idx}: {new_text[:50]} (formatting preserved)")
    
    # Process Slide 2 (Experience)
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]
        print("\n=== SLIDE 2 (Experience) ===")
        
        if len(slide2.shapes) > 1:
            textbox = slide2.shapes[1]
            if hasattr(textbox, "text_frame"):
                tf = textbox.text_frame
                
                # Update company names and roles
                exp = candidate['experience']
                
                # Para 0: Infineon → First company (SOFTLAB360)
                if len(tf.paragraphs) > 0:
                    para = tf.paragraphs[0]
                    update_paragraph_text(para, exp[0]['company'])
                    print(f"  Updated company to: {exp[0]['company']} (formatting preserved)")
                
                # Para 1: Job title and period
                if len(tf.paragraphs) > 1:
                    para = tf.paragraphs[1]
                    update_paragraph_text(para, f"{exp[0]['title']} {exp[0]['period']}")
                    print(f"  Updated title: {exp[0]['title']} (formatting preserved)")
                
                # Para 4+: Replace description
                if len(tf.paragraphs) > 4:
                    para = tf.paragraphs[4]
                    update_paragraph_text(para, f"Area: {exp[0]['area']}")
                    print(f"  Updated area: {exp[0]['area']} (formatting preserved)")
    
    # Save to temporary file first, then copy to final location
    prs.save(temp_output)
    print(f"\n✓ Saved to temporary location")
    
    # Now copy from temp to final (handles OneDrive lock issues)
    import shutil
    import time
    
    try:
        if Path(output_path).exists():
            Path(output_path).unlink()
        time.sleep(0.5)  # Brief pause for file system
        shutil.copy(temp_output, output_path)
        Path(temp_output).unlink()
        print(f"✓ PPTX template populated with formatting preserved: {output_path}")
    except Exception as e:
        print(f"✗ Error moving file: {e}")
        print(f"  Temp file at: {temp_output}")

if __name__ == "__main__":
    populate_template()
